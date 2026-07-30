import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG = RGBColor(248, 250, 252)        # #F8FAFC
    COLOR_DARK_TEXT = RGBColor(15, 23, 42)    # #0F172A
    COLOR_MUTED_TEXT = RGBColor(51, 65, 85)   # #334155
    COLOR_PRIMARY_BLUE = RGBColor(29, 78, 216)# #1D4ED8
    COLOR_GREEN_HEADER = RGBColor(22, 101, 52)# #166534
    COLOR_RED_HEADER = RGBColor(159, 18, 57)  # #9F1239
    
    # Card Backgrounds
    BG_BLUE = RGBColor(239, 246, 255)         # #EFF6FF
    BG_GREEN = RGBColor(240, 253, 244)        # #F0FDF4
    BG_RED = RGBColor(255, 241, 242)          # #FFF1F2
    BG_WHITE = RGBColor(255, 255, 255)

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, subtitle_text, slide_num):
        # Header Box
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.0))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK_TEXT

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = 'Arial'
            p2.font.size = Pt(13)
            p2.font.bold = True
            p2.font.color.rgb = COLOR_PRIMARY_BLUE

        # Slide Number Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.533), Inches(0.4), Inches(1.0), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(224, 231, 255)
        badge.line.fill.background()
        btf = badge.text_frame
        bp = btf.paragraphs[0]
        bp.text = f"{slide_num:02d} / 09"
        bp.alignment = PP_ALIGN.CENTER
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = COLOR_PRIMARY_BLUE

    # -------------------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    card1.fill.solid()
    card1.fill.fore_color.rgb = BG_BLUE
    card1.line.color.rgb = RGBColor(191, 219, 254)

    tf1 = card1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.5)
    tf1.margin_right = Inches(0.5)
    tf1.margin_top = Inches(1.0)

    p1 = tf1.paragraphs[0]
    p1.text = "ENTERPRISE AND VENTURE CREATION (EVC)"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY_BLUE

    p2 = tf1.add_paragraph()
    p2.text = "\nVenture Capitalists (VC)"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_DARK_TEXT

    p3 = tf1.add_paragraph()
    p3.text = "\nFueling Innovation & Scaling High-Growth Startups"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(18)
    p3.font.color.rgb = COLOR_PRIMARY_BLUE

    # -------------------------------------------------------------------------
    # SLIDE 2: Foundations
    # -------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Introduction, Venture Capital & History", "Enterprise and Venture Creation (EVC)", 2)

    # Left Box
    card2_l = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4))
    card2_l.fill.solid()
    card2_l.fill.fore_color.rgb = BG_BLUE
    card2_l.line.color.rgb = RGBColor(191, 219, 254)
    tf2_l = card2_l.text_frame
    tf2_l.word_wrap = True
    tf2_l.margin_left = tf2_l.margin_right = tf2_l.margin_top = tf2_l.margin_bottom = Inches(0.3)

    p = tf2_l.paragraphs[0]
    p.text = "Introduction"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE

    p = tf2_l.add_paragraph()
    p.text = "A startup begins with an idea, but every idea needs money to become a successful business. This is where venture capital plays an important role.\n"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_DARK_TEXT

    p = tf2_l.add_paragraph()
    p.text = "What is Venture Capital?"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE

    bullets_vc = [
      "Venture Capital (VC) is a type of private investment given to startups and small businesses that have high growth potential.",
      "Instead of giving a loan, venture capitalists invest money in exchange for equity, which means ownership in the company.",
      "VC is known as a high-risk, high-return investment because many startups fail, but successful ones can generate huge profits."
    ]
    for b in bullets_vc:
        p = tf2_l.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_DARK_TEXT

    # Right Box
    card2_r = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.5), Inches(5.7), Inches(5.4))
    card2_r.fill.solid()
    card2_r.fill.fore_color.rgb = BG_RED
    card2_r.line.color.rgb = RGBColor(254, 205, 211)
    tf2_r = card2_r.text_frame
    tf2_r.word_wrap = True
    tf2_r.margin_left = tf2_r.margin_right = tf2_r.margin_top = tf2_r.margin_bottom = Inches(0.3)

    p = tf2_r.paragraphs[0]
    p.text = "Who is a Venture Capitalist?"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED_HEADER

    p = tf2_r.add_paragraph()
    p.text = "A Venture Capitalist is an investor or investment firm that provides funding to startups."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK_TEXT

    p = tf2_r.add_paragraph()
    p.text = "Apart from money, they also provide: Business guidance, Industry connections, Marketing support, Financial advice, and Mentorship.\n"
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_DARK_TEXT

    p = tf2_r.add_paragraph()
    p.text = "History of Venture Capital"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED_HEADER

    bullets_hist = [
      "Before 1946, entrepreneurs mainly depended on family, friends, or banks.",
      "In 1946, the first modern VC firm, ARDC, was founded by Georges Doriot (Father of Venture Capital).",
      "Its investment in Digital Equipment Corporation proved startups generate massive returns.",
      "Today, Google, Apple, Amazon, and Meta all scaled via early venture capital."
    ]
    for b in bullets_hist:
        p = tf2_r.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK_TEXT

    # -------------------------------------------------------------------------
    # SLIDE 3: Need, Features & Objectives
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "Need, Features & Objectives", "Enterprise and Venture Creation (EVC)", 3)

    cols = [
      ("Why is Venture Capital Needed?", BG_RED, RGBColor(254, 205, 211), COLOR_RED_HEADER, [
        "Startups face 3 major problems:",
        "• Lack of money",
        "• Lack of business experience",
        "• Lack of professional contacts",
        "",
        "Solution: Venture capital solves these problems by providing funding, mentorship, and networking opportunities."
      ]),
      ("Features of Venture Capital", BG_BLUE, RGBColor(191, 219, 254), COLOR_PRIMARY_BLUE, [
        "• Equity-based investment",
        "• High risk and high return",
        "• Long-term investment horizon",
        "• Active involvement in business decisions",
        "• Focus on innovation and technology",
        "• Supports scalable high-growth models"
      ]),
      ("Objectives of Venture Capital", BG_GREEN, RGBColor(187, 247, 208), COLOR_GREEN_HEADER, [
        "• Promote entrepreneurship",
        "• Encourage tech innovation",
        "• Generate employment opportunities",
        "• Support early-stage startup growth",
        "• Increase overall economic development",
        "• Earn high returns through successful exits"
      ])
    ]

    for i, (col_title, bg_c, border_c, head_c, lines) in enumerate(cols):
        left_pos = Inches(0.8 + i * 3.95)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.5), Inches(3.8), Inches(5.4))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_c
        card.line.color.rgb = border_c
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = col_title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = head_c

        for line in lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(11.5)
            p.font.color.rgb = COLOR_DARK_TEXT

    # -------------------------------------------------------------------------
    # SLIDE 4: Venture Capital Investment Process
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Venture Capital Investment Process", "The 7 Steps in Venture Capital Investment", 4)

    steps_data = [
      ("Step 1: Business Idea", "Entrepreneur develops an innovative business idea."),
      ("Step 2: Business Plan", "Prepares market analysis, financial projections & goals."),
      ("Step 3: Pitching", "Entrepreneur presents idea to venture capitalists."),
      ("Step 4: Due Diligence", "VC examines business model, team, market & legal docs."),
      ("Step 5: Investment", "VC invests funds in exchange for equity ownership."),
      ("Step 6: Growth", "Startup uses funds for R&D, hiring & market expansion."),
      ("Step 7: Exit", "VC earns profits via IPO, Acquisition, or Share Sale.")
    ]

    for i, (st_title, st_desc) in enumerate(steps_data):
        row = i // 4
        col = i % 4
        left_p = Inches(0.8 + col * 2.95)
        top_p = Inches(1.5 + row * 2.7)

        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_p, top_p, Inches(2.8), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_BLUE if (i % 2 == 0) else BG_GREEN
        card.line.color.rgb = RGBColor(191, 219, 254) if (i % 2 == 0) else RGBColor(187, 247, 208)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = st_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_BLUE if (i % 2 == 0) else COLOR_GREEN_HEADER

        p = tf.add_paragraph()
        p.text = st_desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK_TEXT

    # -------------------------------------------------------------------------
    # SLIDE 5: Funding Stages & Comparisons
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Funding Stages & Comparisons", "Funding Stages & Venture Capital vs Angel Investor", 5)

    # Left: Funding Stages
    card5_l = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4))
    card5_l.fill.solid()
    card5_l.fill.fore_color.rgb = BG_BLUE
    card5_l.line.color.rgb = RGBColor(191, 219, 254)
    tf5_l = card5_l.text_frame
    tf5_l.word_wrap = True
    tf5_l.margin_left = tf5_l.margin_right = tf5_l.margin_top = tf5_l.margin_bottom = Inches(0.3)

    p = tf5_l.paragraphs[0]
    p.text = "Funding Stages"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE

    stages = [
      ("Pre-Seed", "Idea, research, and proof of concept stage."),
      ("Seed Funding", "Product development and prototype creation."),
      ("Series A", "Business expansion and customer acquisition."),
      ("Series B", "Scaling operations and entering new markets."),
      ("Series C", "International expansion and major acquisitions."),
      ("IPO", "Initial Public Offering - company becomes publicly listed.")
    ]
    for s_name, s_desc in stages:
        p = tf5_l.add_paragraph()
        p.text = f"• {s_name}: {s_desc}"
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_DARK_TEXT

    # Right: Angel vs VC
    card5_r = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.5), Inches(5.7), Inches(5.4))
    card5_r.fill.solid()
    card5_r.fill.fore_color.rgb = BG_WHITE
    card5_r.line.color.rgb = RGBColor(226, 232, 240)
    tf5_r = card5_r.text_frame
    tf5_r.word_wrap = True
    tf5_r.margin_left = tf5_r.margin_right = tf5_r.margin_top = tf5_r.margin_bottom = Inches(0.3)

    p = tf5_r.paragraphs[0]
    p.text = "Venture Capital vs Angel Investor"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK_TEXT

    p = tf5_r.add_paragraph()
    p.text = "Angel Investor:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED_HEADER
    for a in ["Individual wealthy investor", "Invests personal money", "Smaller investment size", "Invests at very early/idea stage"]:
        p = tf5_r.add_paragraph()
        p.text = f"  • {a}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK_TEXT

    p = tf5_r.add_paragraph()
    p.text = "\nVenture Capitalist:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN_HEADER
    for v in ["Professional investment firm", "Invests pooled funds from institutional LPs", "Larger investment ticket size", "Invests at high-growth & scaling stage"]:
        p = tf5_r.add_paragraph()
        p.text = f"  • {v}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK_TEXT

    # -------------------------------------------------------------------------
    # SLIDE 6: VC vs Bank Loan, Risks & Startup Failure (DETAILED EXPLANATIONS)
    # -------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "VC vs Bank Loan, Risks & Startup Failure", "Enterprise and Venture Creation (EVC)", 6)

    # Top Box: VC vs Bank Loan
    card6_top = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(2.5))
    card6_top.fill.solid()
    card6_top.fill.fore_color.rgb = BG_WHITE
    card6_top.line.color.rgb = RGBColor(226, 232, 240)
    tf6_t = card6_top.text_frame
    tf6_t.word_wrap = True
    tf6_t.margin_left = tf6_t.margin_right = tf6_t.margin_top = tf6_t.margin_bottom = Inches(0.2)

    p = tf6_t.paragraphs[0]
    p.text = "1. Venture Capital vs. Bank Loan Comparison (Detailed Breakdown)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE

    p = tf6_t.add_paragraph()
    p.text = "Venture Capital:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    vc_details = [
      ("Equity investment", "Capital given in exchange for company ownership share rather than debt."),
      ("No EMI required", "No fixed monthly repayments or interest obligations regardless of profit."),
      ("Ownership shared", "Founders surrender equity stake and share company control with investor partners."),
      ("Strategic support", "Active business guidance, industry connections, and mentorship provided.")
    ]
    for title, desc in vc_details:
        p = tf6_t.add_paragraph()
        p.text = f"  • {title}: {desc}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_DARK_TEXT

    p = tf6_t.add_paragraph()
    p.text = "Bank Loan:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED_HEADER
    bank_details = [
      ("Debt financing", "Borrowed money with strict interest rates and full principal repayment terms."),
      ("Mandatory EMI", "Regular monthly payments required regardless of business profits or losses."),
      ("100% ownership kept", "Founder retains full company ownership, equity, and independent decision authority."),
      ("No business support", "Banks provide money only; no strategic advice, networking, or mentorship.")
    ]
    for title, desc in bank_details:
        p = tf6_t.add_paragraph()
        p.text = f"  • {title}: {desc}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_DARK_TEXT

    # Bottom Left: Risks
    card6_bl = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.0), Inches(5.7), Inches(3.0))
    card6_bl.fill.solid()
    card6_bl.fill.fore_color.rgb = BG_RED
    card6_bl.line.color.rgb = RGBColor(254, 205, 211)
    tf6_bl = card6_bl.text_frame
    tf6_bl.word_wrap = True
    tf6_bl.margin_left = tf6_bl.margin_right = tf6_bl.margin_top = tf6_bl.margin_bottom = Inches(0.2)

    p = tf6_bl.paragraphs[0]
    p.text = "2. Investment & Entrepreneur Risks"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED_HEADER

    p = tf6_bl.add_paragraph()
    p.text = "For Investors:"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED_HEADER
    inv_risks = [
      ("Startup Failure Risk", "Over 90% of early-stage startups fail, resulting in total loss of invested capital."),
      ("Market & Tech Competition", "Rival companies or fast technological changes can render the business obsolete."),
      ("Illiquidity & Long Gestation", "Capital is locked in private shares for 7–10 years before any exit opportunity.")
    ]
    for title, desc in inv_risks:
        p = tf6_bl.add_paragraph()
        p.text = f"  • {title}: {desc}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_DARK_TEXT

    p = tf6_bl.add_paragraph()
    p.text = "\nFor Entrepreneurs:"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED_HEADER
    ent_risks = [
      ("Loss of Ownership & Dilution", "Selling equity dilutes founders' control and equity stake in the company."),
      ("High Investor Growth Pressure", "VCs demand 10x-100x financial returns, pushing founders toward aggressive targets."),
      ("Board Control & Conflicts", "Investors secure board seats and veto rights, restricting founder autonomy.")
    ]
    for title, desc in ent_risks:
        p = tf6_bl.add_paragraph()
        p.text = f"  • {title}: {desc}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_DARK_TEXT

    # Bottom Right: Failure
    card6_br = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(4.0), Inches(5.7), Inches(3.0))
    card6_br.fill.solid()
    card6_br.fill.fore_color.rgb = BG_RED
    card6_br.line.color.rgb = RGBColor(254, 205, 211)
    tf6_br = card6_br.text_frame
    tf6_br.word_wrap = True
    tf6_br.margin_left = tf6_br.margin_right = tf6_br.margin_top = tf6_br.margin_bottom = Inches(0.2)

    p = tf6_br.paragraphs[0]
    p.text = "3. What Happens If a Startup Fails?"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED_HEADER

    p = tf6_br.add_paragraph()
    p.text = "Key consequences when a VC-backed startup fails:"
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED_HEADER

    fail_points = [
      ("VC Absorbs Capital Loss", "The venture capital firm loses its entire equity investment without collateral recourse."),
      ("No Personal Liability for Founders", "Founders do not repay equity money because equity carries risk, unlike bank debt."),
      ("Workforce Layoffs & Closure", "Operations cease, leading to employee layoffs and workplace shutdown."),
      ("Asset Liquidation & Distribution", "Remaining physical assets and patents are sold off to satisfy debt creditors.")
    ]
    for title, desc in fail_points:
        p = tf6_br.add_paragraph()
        p.text = f"  • {title}: {desc}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_DARK_TEXT

    # -------------------------------------------------------------------------
    # SLIDE 7: Venture Capital in India (DETAILED EXPLANATIONS)
    # -------------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Venture Capital in India", "India is one of the fastest-growing startup ecosystems worldwide.", 7)

    # Left: Govt & Sectors
    card7_l = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4))
    card7_l.fill.solid()
    card7_l.fill.fore_color.rgb = BG_GREEN
    card7_l.line.color.rgb = RGBColor(187, 247, 208)
    tf7_l = card7_l.text_frame
    tf7_l.word_wrap = True
    tf7_l.margin_left = tf7_l.margin_right = tf7_l.margin_top = tf7_l.margin_bottom = Inches(0.2)

    p = tf7_l.paragraphs[0]
    p.text = "Government Initiatives"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN_HEADER

    govt_items = [
      ("Startup India", "Provides tax exemptions, fast-track patent filings, and simplified compliance for startups."),
      ("Digital India", "Expands nationwide high-speed internet and digital infrastructure to support tech startups."),
      ("Atal Innovation Mission", "Establishes world-class incubators, tinkering labs, and innovation hubs across educational institutions."),
      ("Fund of Funds (FFS)", "₹10,000 Crore government corpus managed by SIDBI to inject capital into registered VCs.")
    ]
    for name, desc in govt_items:
        p = tf7_l.add_paragraph()
        p.text = f"• {name}: {desc}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_DARK_TEXT

    p = tf7_l.add_paragraph()
    p.text = "\nMajor Investment Sectors"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN_HEADER

    sector_items = [
      ("Artificial Intelligence (AI)", "Building generative models, predictive analytics, and automated enterprise tools."),
      ("FinTech", "Revolutionizing digital payments via UPI, neo-banking, and micro-lending platforms."),
      ("HealthTech", "Providing online telemedicine, AI diagnostic tools, and accessible digital healthcare."),
      ("EdTech", "Enabling interactive online education, skill upskilling platforms, and exam preparation."),
      ("AgriTech", "Optimizing supply chain logistics, IoT soil sensors, and direct farm-to-consumer platforms."),
      ("Electric Vehicles (EV)", "Building battery swapping networks, EV manufacturing, and clean green mobility."),
      ("SaaS (Software as a Service)", "Developing cloud-based enterprise software built in India for global customers.")
    ]
    for name, desc in sector_items:
        p = tf7_l.add_paragraph()
        p.text = f"• {name}: {desc}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_DARK_TEXT

    # Right: VCs & Startups
    card7_r = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.5), Inches(5.7), Inches(5.4))
    card7_r.fill.solid()
    card7_r.fill.fore_color.rgb = BG_BLUE
    card7_r.line.color.rgb = RGBColor(191, 219, 254)
    tf7_r = card7_r.text_frame
    tf7_r.word_wrap = True
    tf7_r.margin_left = tf7_r.margin_right = tf7_r.margin_top = tf7_r.margin_bottom = Inches(0.2)

    p = tf7_r.paragraphs[0]
    p.text = "Top Venture Capital Firms in India"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE

    vc_firms = [
      ("Peak XV Partners", "Formerly Sequoia India; leading multi-stage VC firm backing iconic tech unicorns."),
      ("Accel India", "Renowned early-stage investor behind global giants like Flipkart, Swiggy, and Freshworks."),
      ("Kalaari Capital", "Early-stage fund focusing on consumer internet, e-commerce, and deep tech innovations."),
      ("Blume Ventures", "Prominent seed and pre-Series A fund empowering home-grown Indian tech startups."),
      ("Nexus Venture Partners", "Pioneering Indo-US venture firm supporting enterprise software, AI, and commercial tech.")
    ]
    for name, desc in vc_firms:
        p = tf7_r.add_paragraph()
        p.text = f"• {name}: {desc}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_DARK_TEXT

    p = tf7_r.add_paragraph()
    p.text = "\nFamous VC-Funded Indian Startups"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE

    famous_s = [
      ("Flipkart", "E-commerce pioneer backed by Accel & SoftBank, acquired by Walmart for $16 Billion."),
      ("Ola", "Mobility and EV leader transforming urban ride-hailing and clean transit across India."),
      ("Swiggy", "On-demand food delivery and quick-commerce innovator with Instamart grocery delivery."),
      ("Zomato", "Food delivery & restaurant discovery marketplace successfully listed on public stock exchanges."),
      ("Meesho", "Social commerce platform democratizing online selling for millions of small businesses.")
    ]
    for name, desc in famous_s:
        p = tf7_r.add_paragraph()
        p.text = f"• {name}: {desc}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_DARK_TEXT

    # -------------------------------------------------------------------------
    # SLIDE 8: Real-Life Examples
    # -------------------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Real-Life Examples: From Startup to Major Brand", "How VC Funding Transformed Small Ideas into Household Brands", 8)

    card8_l = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4))
    card8_l.fill.solid()
    card8_l.fill.fore_color.rgb = BG_BLUE
    card8_l.line.color.rgb = RGBColor(191, 219, 254)
    tf8_l = card8_l.text_frame
    tf8_l.word_wrap = True
    tf8_l.margin_left = tf8_l.margin_right = tf8_l.margin_top = tf8_l.margin_bottom = Inches(0.3)

    p = tf8_l.paragraphs[0]
    p.text = "1. Real-Life Conceptual Example"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE

    farm_points = [
      "Suppose 4 engineering students develop an AI-based farming application.",
      "They need ₹50 Lakhs to build and launch the product.",
      "A venture capital firm invests ₹50 Lakhs in exchange for 20% equity.",
      "The startup grows successfully over 5 years and reaches a valuation of ₹100 Crores.",
      "The VC's 20% equity is now worth ₹20 Crores, while founders retain majority ownership of a massive enterprise!"
    ]
    for pt in farm_points:
        p = tf8_l.add_paragraph()
        p.text = f"• {pt}\n"
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_DARK_TEXT

    card8_r = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.5), Inches(5.7), Inches(5.4))
    card8_r.fill.solid()
    card8_r.fill.fore_color.rgb = BG_GREEN
    card8_r.line.color.rgb = RGBColor(187, 247, 208)
    tf8_r = card8_r.text_frame
    tf8_r.word_wrap = True
    tf8_r.margin_left = tf8_r.margin_right = tf8_r.margin_top = tf8_r.margin_bottom = Inches(0.3)

    p = tf8_r.paragraphs[0]
    p.text = "2. Flipkart (Real-Time Major Brand Case)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN_HEADER

    fk_points = [
      "Initial Startup (2007): Started in a 2BHK apartment in Bengaluru by Sachin & Binny Bansal as an online bookstore with ₹4 Lakhs savings.",
      "VC Funding Boost (2009): Accel India invested $1 Million (₹4.5 Crores) when traditional banks refused loans. Later, Tiger Global & SoftBank backed them.",
      "Growth to Major Brand: Built nationwide logistics, Cash-on-Delivery, and expanded across e-commerce categories.",
      "Landmark Exit (2018): Walmart acquired a controlling stake for $16 Billion (₹1.1 Lakh Crores) — giving early VC Accel a 300x Return on Investment!"
    ]
    for pt in fk_points:
        p = tf8_r.add_paragraph()
        p.text = f"★ {pt}\n"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK_TEXT

    # -------------------------------------------------------------------------
    # SLIDE 9: Advantages, Disadvantages & Conclusion (DETAILED EXPLANATIONS)
    # -------------------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "Advantages, Disadvantages & Conclusion", "Enterprise and Venture Creation (EVC)", 9)

    # Top Left: Advantages
    card9_tl = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(3.6))
    card9_tl.fill.solid()
    card9_tl.fill.fore_color.rgb = BG_GREEN
    card9_tl.line.color.rgb = RGBColor(187, 247, 208)
    tf9_tl = card9_tl.text_frame
    tf9_tl.word_wrap = True
    tf9_tl.margin_left = tf9_tl.margin_right = tf9_tl.margin_top = tf9_tl.margin_bottom = Inches(0.2)

    p = tf9_tl.paragraphs[0]
    p.text = "Advantages of Venture Capital"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN_HEADER

    adv_items = [
      ("Substantial Capital Access", "Provides large-scale growth funding without monthly interest payments or collateral debt."),
      ("Strategic Mentorship", "VC partners bring experienced business guidance, board leadership, and operational advice."),
      ("Accelerated Scaling", "Enables rapid market expansion, aggressive marketing campaigns, and top-tier hiring."),
      ("Strong Industry Connections", "Unlocks valuable corporate partnerships, key enterprise clients, and future funding networks."),
      ("Increased Credibility & Trust", "VC backing serves as a strong stamp of validation, boosting market reputation and trust.")
    ]
    for title, desc in adv_items:
        p = tf9_tl.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_DARK_TEXT

    # Top Right: Disadvantages
    card9_tr = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.4), Inches(5.7), Inches(3.6))
    card9_tr.fill.solid()
    card9_tr.fill.fore_color.rgb = BG_RED
    card9_tr.line.color.rgb = RGBColor(254, 205, 211)
    tf9_tr = card9_tr.text_frame
    tf9_tr.word_wrap = True
    tf9_tr.margin_left = tf9_tr.margin_right = tf9_tr.margin_top = tf9_tr.margin_bottom = Inches(0.2)

    p = tf9_tr.paragraphs[0]
    p.text = "Disadvantages of Venture Capital"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED_HEADER

    disadv_items = [
      ("Equity & Ownership Loss", "Founders must give up significant equity ownership shares and dilute future profits."),
      ("High Investor Return Expectations", "VCs demand 10x-100x financial returns, placing immense performance pressure on founders."),
      ("Difficult & Lengthy Process", "Pitching, due diligence, and deal negotiation can take 6-12 months of intense effort."),
      ("Pressure to Grow Quickly", "Forces startups to prioritize aggressive hyper-scaling over organic development.")
    ]
    for title, desc in disadv_items:
        p = tf9_tr.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_DARK_TEXT

    # Bottom Box: Conclusion
    card9_b = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.1), Inches(11.733), Inches(1.9))
    card9_b.fill.solid()
    card9_b.fill.fore_color.rgb = BG_BLUE
    card9_b.line.color.rgb = RGBColor(191, 219, 254)
    tf9_b = card9_b.text_frame
    tf9_b.word_wrap = True
    tf9_b.margin_left = tf9_b.margin_right = tf9_b.margin_top = tf9_b.margin_bottom = Inches(0.2)

    p = tf9_b.paragraphs[0]
    p.text = "Conclusion"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_BLUE

    p = tf9_b.add_paragraph()
    p.text = "To conclude, venture capital is an indispensable growth engine for innovative, high-potential startups. Beyond capital funding, venture capitalists provide vital mentorship, strategic networking, and operational expertise. Although it involves significant risk, equity dilution, and intense return expectations, venture capital has built iconic global and Indian companies—driving technology innovation, employment generation, and national economic growth."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_DARK_TEXT

    p = tf9_b.add_paragraph()
    p.text = "\nThank You!"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN_HEADER

    # Save presentation
    prs.save(output_path)
    print(f"Successfully generated PowerPoint presentation at: {output_path}")

if __name__ == "__main__":
    output_desktop = os.path.expanduser("~/Desktop/Venture_Capital_Presentation.pptx")
    create_presentation(output_desktop)
